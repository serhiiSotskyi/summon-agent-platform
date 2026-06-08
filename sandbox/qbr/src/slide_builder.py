from __future__ import annotations

from copy import deepcopy
from pathlib import Path
from typing import Iterable, Sequence

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


class SlideBuilder:
    def __init__(self, template_path: str | Path, chart_styles: dict | None = None) -> None:
        self.template_path = Path(template_path)
        self.template_prs = Presentation(str(self.template_path))
        self.prs = self._create_output_presentation(self.template_prs)

        colors = (chart_styles or {}).get("colors", {})
        fonts = (chart_styles or {}).get("fonts", {})
        self.text_primary = _hex_to_rgb(colors.get("text_primary", "#142A4D"))
        self.text_secondary = _hex_to_rgb(colors.get("text_secondary", "#5A5A5A"))
        self.text_body = _hex_to_rgb(colors.get("text_body", "#2D2D2D"))
        self.surface = _hex_to_rgb(colors.get("surface", "#FFFFFF"))
        self.border = _hex_to_rgb(colors.get("border", "#D9D9D9"))
        self.accent = _hex_to_rgb(colors.get("accent", "#C32026"))
        self.positive = _hex_to_rgb(colors.get("positive", "#111111"))
        self.negative = _hex_to_rgb(colors.get("negative", "#C32026"))
        self.table_header = _hex_to_rgb(colors.get("table_header", "#DBE6F4"))
        self.table_total = _hex_to_rgb(colors.get("table_total", "#EBF1FA"))
        self.title_size = int(fonts.get("title_size", 16))
        self.body_size = int(fonts.get("body_size", 11))

    def add_title_slide(self, title: str, subtitle: str) -> None:
        slide = self._get_template_slide("TITLE")
        if slide is None:
            slide = self._new_slide()
            self._add_title(slide, title, top=0.9, font_size=38)
            self._add_subtitle(slide, subtitle, top=1.8, font_size=22)
            return

        self._replace_text_placeholder(slide, "{{TITLE}}", title)
        self._replace_text_placeholder(slide, "{{SUBTITLE}}", subtitle)
        self._clear_text_placeholder(slide, "{{SLIDE_TYPE:TITLE}}")

    def add_divider_slide(self, title: str) -> None:
        slide = self._get_template_slide("DIVIDER")
        if slide is None:
            slide = self._new_slide()
            self._add_title(slide, title, top=2.8, font_size=56)
            return

        self._replace_text_placeholder(slide, "{{TITLE}}", title)
        self._clear_text_placeholder(slide, "{{SLIDE_TYPE:DIVIDER}}")

    def add_trend_slide(
        self,
        title: str,
        subtitle: str,
        cpl_cvr_chart_path: str | Path,
        cost_leads_chart_path: str | Path,
        bullets: Iterable[str] | None = None,
        use_template: bool = True,
    ) -> None:
        slide = self._get_template_slide("TREND") if use_template else None
        if slide is None:
            slide = self._new_slide()
            self._add_title(slide, title)
            self._add_subtitle(slide, subtitle)
            slide.shapes.add_picture(str(cost_leads_chart_path), Inches(0.7), Inches(1.55), width=Inches(5.45))
            slide.shapes.add_picture(str(cpl_cvr_chart_path), Inches(6.7), Inches(1.55), width=Inches(5.2))
            if bullets:
                self._add_bullets(slide, bullets, left=0.8, top=4.85, width=12.0, height=1.55)
            return

        self._replace_text_placeholder(slide, "{{TITLE}}", title)
        self._replace_text_placeholder(slide, "{{SUBTITLE}}", subtitle)
        self._replace_picture_placeholder(slide, "{{CHART_LEFT}}", cost_leads_chart_path)
        self._replace_picture_placeholder(slide, "{{CHART_RIGHT}}", cpl_cvr_chart_path)
        self._replace_bullets_placeholder(slide, "{{BULLETS}}", bullets or [])
        self._clear_text_placeholder(slide, "{{SLIDE_TYPE:TREND}}")

    def add_mix_slide(
        self,
        title: str,
        subtitle: str,
        cost_mix_chart_path: str | Path,
        leads_mix_chart_path: str | Path,
        bullets: Iterable[str],
    ) -> None:
        slide = self._get_template_slide("MIX")
        if slide is None:
            slide = self._new_slide()
            self._add_title(slide, title)
            self._add_subtitle(slide, subtitle)
            slide.shapes.add_picture(str(cost_mix_chart_path), Inches(0.9), Inches(1.45), width=Inches(5.6))
            slide.shapes.add_picture(str(leads_mix_chart_path), Inches(6.9), Inches(1.45), width=Inches(5.6))
            self._add_bullets(slide, bullets, left=0.8, top=5.25, width=12.0, height=1.7)
            return

        self._replace_text_placeholder(slide, "{{TITLE}}", title)
        self._replace_text_placeholder(slide, "{{SUBTITLE}}", subtitle)
        self._replace_picture_placeholder(slide, "{{CHART_LEFT}}", cost_mix_chart_path)
        self._replace_picture_placeholder(slide, "{{CHART_RIGHT}}", leads_mix_chart_path)
        self._replace_bullets_placeholder(slide, "{{BULLETS}}", bullets)
        self._clear_text_placeholder(slide, "{{SLIDE_TYPE:MIX}}")

    def add_table_slide(
        self,
        title: str,
        subtitle: str,
        table_df: pd.DataFrame,
        bullets: Iterable[str],
    ) -> None:
        slide = self._get_template_slide("TABLE")
        if slide is None:
            slide = self._new_slide()
            self._add_title(slide, title)
            self._add_subtitle(slide, subtitle)
            self._render_table(slide, table_df, Inches(0.4), Inches(1.35), Inches(12.5), Inches(3.5))
            self._add_bullets(slide, bullets, left=0.8, top=5.05, width=12.0, height=2.0)
            return

        self._replace_text_placeholder(slide, "{{TITLE}}", title)
        self._replace_text_placeholder(slide, "{{SUBTITLE}}", subtitle)
        self._replace_table_placeholder(slide, "{{TABLE_DATA}}", table_df)
        self._replace_bullets_placeholder(slide, "{{BULLETS}}", bullets)
        self._clear_text_placeholder(slide, "{{SLIDE_TYPE:TABLE}}")

    def add_summary_cards_slide(
        self,
        title: str,
        subtitle: str,
        kpis: Sequence[dict],
        bullets: Iterable[str] | None = None,
    ) -> None:
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_subtitle(slide, subtitle)
        self._render_kpi_cards(slide, kpis, left=0.65, top=1.45, width=12.0, height=3.4)
        if bullets:
            self._add_bullets(slide, bullets, left=0.8, top=5.2, width=11.8, height=1.35)

    def add_single_chart_slide(
        self,
        title: str,
        subtitle: str,
        chart_path: str | Path,
        bullets: Iterable[str],
        source_note: str = "",
    ) -> None:
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_subtitle(slide, subtitle)
        slide.shapes.add_picture(str(chart_path), Inches(0.55), Inches(1.35), width=Inches(7.35), height=Inches(4.6))
        self._add_bullets(slide, bullets, left=8.15, top=1.55, width=4.3, height=3.8)
        if source_note:
            self._add_source_note(slide, source_note)

    def add_auction_insights_slide(
        self,
        title: str,
        subtitle: str,
        table_df: pd.DataFrame,
        bullets: Iterable[str],
        source_note: str = "",
    ) -> None:
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_subtitle(slide, subtitle)
        self._render_table(slide, table_df.head(8), Inches(0.45), Inches(1.35), Inches(12.35), Inches(3.35))
        self._add_bullets(slide, bullets, left=0.8, top=4.95, width=11.7, height=1.55)
        if source_note:
            self._add_source_note(slide, source_note)

    def add_recommendations_slide(
        self,
        title: str,
        subtitle: str,
        recommendations: Sequence[dict],
        source_note: str = "",
    ) -> None:
        slide = self._new_slide()
        self._add_title(slide, title)
        self._add_subtitle(slide, subtitle)
        self._add_structured_recommendations(slide, recommendations, left=0.8, top=1.55, width=11.8, height=4.8)
        if source_note:
            self._add_source_note(slide, source_note)

    def save(self, output_path: str | Path) -> None:
        output = Path(output_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output))

    def _get_template_slide(self, slide_kind: str):
        for template_slide in self.template_prs.slides:
            if self._slide_contains_text(template_slide, f"{{{{SLIDE_TYPE:{slide_kind}}}}}"):
                return self._clone_template_slide(template_slide)
        return None

    def _clone_template_slide(self, template_slide):
        slide = self._new_slide()
        for shape in template_slide.shapes:
            slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")
        return slide

    def _new_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(247, 249, 252)
        return slide

    def _add_title(self, slide, text: str, top: float = 0.28, font_size: int = 30) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.3), Inches(0.8)).text_frame
        box.clear()
        para = box.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = self.text_primary

    def _add_subtitle(self, slide, text: str, top: float = 0.94, font_size: int = 15) -> None:
        box = slide.shapes.add_textbox(Inches(0.6), Inches(top), Inches(12.3), Inches(0.45)).text_frame
        para = box.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = self.text_secondary

    def _add_bullets(
        self,
        slide,
        bullets: Iterable[str],
        left: float,
        top: float,
        width: float,
        height: float,
    ) -> None:
        text_frame = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
        self._populate_bullets(text_frame, bullets)

    def _add_source_note(self, slide, text: str) -> None:
        text_frame = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(6.0), Inches(0.3)).text_frame
        para = text_frame.paragraphs[0]
        run = para.add_run()
        run.text = text
        run.font.size = Pt(9)
        run.font.color.rgb = self.text_secondary

    def _add_structured_recommendations(self, slide, recommendations: Sequence[dict], left: float, top: float, width: float, height: float) -> None:
        text_frame = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height)).text_frame
        text_frame.clear()
        recs = list(recommendations) or [{"heading": "Next quarter focus", "text": "No recommendation data was available."}]
        for idx, item in enumerate(recs):
            heading = str(item.get("heading", "")).strip()
            text = str(item.get("text", "")).strip()
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            if heading:
                run = paragraph.add_run()
                run.text = f"{heading}: "
                run.font.bold = True
                run.font.size = Pt(16)
                run.font.color.rgb = self.text_primary
            run = paragraph.add_run()
            run.text = text
            run.font.size = Pt(16)
            run.font.color.rgb = self.text_body
            paragraph.level = 0
            paragraph.space_after = Pt(12)

    def _render_kpi_cards(self, slide, kpis: Sequence[dict], left: float, top: float, width: float, height: float) -> None:
        cards = list(kpis)
        if not cards:
            return

        columns = 4
        rows = 2
        horizontal_gap = 0.18
        vertical_gap = 0.18
        card_width = (width - horizontal_gap * (columns - 1)) / columns
        card_height = (height - vertical_gap * (rows - 1)) / rows

        for index, kpi in enumerate(cards[: columns * rows]):
            row = index // columns
            column = index % columns
            card_left = left + column * (card_width + horizontal_gap)
            card_top = top + row * (card_height + vertical_gap)

            card = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(card_left),
                Inches(card_top),
                Inches(card_width),
                Inches(card_height),
            )
            card.fill.solid()
            card.fill.fore_color.rgb = self.surface
            card.line.color.rgb = self.border
            card.line.width = Pt(1)

            accent = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                Inches(card_left),
                Inches(card_top),
                Inches(card_width),
                Inches(0.06),
            )
            accent.fill.solid()
            accent.fill.fore_color.rgb = self.accent
            accent.line.color.rgb = self.accent

            label_box = slide.shapes.add_textbox(
                Inches(card_left + 0.14),
                Inches(card_top + 0.16),
                Inches(card_width - 0.28),
                Inches(0.28),
            ).text_frame
            label_para = label_box.paragraphs[0]
            label_run = label_para.add_run()
            label_run.text = str(kpi.get("label", "Metric"))
            label_run.font.size = Pt(10)
            label_run.font.bold = True
            label_run.font.color.rgb = self.text_secondary

            value_box = slide.shapes.add_textbox(
                Inches(card_left + 0.14),
                Inches(card_top + 0.48),
                Inches(card_width - 0.28),
                Inches(0.42),
            ).text_frame
            value_para = value_box.paragraphs[0]
            value_run = value_para.add_run()
            value_run.text = str(kpi.get("value", "n/a"))
            value_run.font.size = Pt(20)
            value_run.font.bold = True
            value_run.font.color.rgb = self.text_primary

            yoy_box = slide.shapes.add_textbox(
                Inches(card_left + 0.14),
                Inches(card_top + card_height - 0.44),
                Inches(card_width - 0.28),
                Inches(0.24),
            ).text_frame
            yoy_para = yoy_box.paragraphs[0]
            yoy_run = yoy_para.add_run()
            yoy_run.text = f"YoY: {kpi.get('yoy_label', 'n/a')}"
            yoy_run.font.size = Pt(10)
            yoy_run.font.bold = True
            yoy_run.font.color.rgb = self._resolve_yoy_color(str(kpi.get("key", "")), kpi.get("yoy"))

    def _replace_text_placeholder(self, slide, placeholder: str, value: str) -> bool:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, value)
                        return True
            if placeholder == shape.text_frame.text:
                shape.text_frame.text = value
                return True
        return False

    def _clear_text_placeholder(self, slide, placeholder: str) -> None:
        self._replace_text_placeholder(slide, placeholder, "")

    def _replace_bullets_placeholder(self, slide, placeholder: str, bullets: Iterable[str]) -> bool:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            if shape.text_frame.text.strip() == placeholder:
                self._populate_bullets(shape.text_frame, bullets)
                return True
        return False

    def _replace_picture_placeholder(self, slide, placeholder: str, image_path: str | Path) -> bool:
        target = self._find_placeholder_shape(slide, placeholder)
        if target is None:
            return False

        left, top, width, height = target.left, target.top, target.width, target.height
        self._remove_shape(target)
        slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
        return True

    def _replace_table_placeholder(self, slide, placeholder: str, table_df: pd.DataFrame) -> bool:
        target = self._find_placeholder_shape(slide, placeholder)
        if target is None:
            return False

        left, top, width, height = target.left, target.top, target.width, target.height
        self._remove_shape(target)
        self._render_table(slide, table_df, left, top, width, height)
        return True

    def _render_table(self, slide, table_df: pd.DataFrame, left, top, width, height) -> None:
        safe_table_df = table_df.copy()
        if safe_table_df.empty:
            safe_table_df = pd.DataFrame([{"Status": "No data available"}])

        n_rows = len(safe_table_df) + 1
        n_cols = len(safe_table_df.columns)
        table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = table_shape.table

        for col_idx, col_name in enumerate(safe_table_df.columns):
            cell = table.cell(0, col_idx)
            cell.text = str(col_name)
            self._style_cell_text(cell, bold=True, size=10)
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.table_header

        for row_idx, (_, row) in enumerate(safe_table_df.iterrows(), start=1):
            is_total = str(row.iloc[0]).strip().lower() == "total"
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(value)
                self._style_cell_text(cell, bold=is_total, size=10)
                if is_total:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = self.table_total

    def _populate_bullets(self, text_frame, bullets: Iterable[str]) -> None:
        text_frame.clear()
        bullets_list = [str(bullet) for bullet in bullets if str(bullet).strip()]
        if not bullets_list:
            return
        for idx, bullet in enumerate(bullets_list):
            paragraph = text_frame.paragraphs[0] if idx == 0 else text_frame.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = self.text_body

    def _find_placeholder_shape(self, slide, placeholder: str):
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.text.strip() == placeholder:
                return shape
        return None

    @staticmethod
    def _slide_contains_text(slide, text: str) -> bool:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and text in shape.text_frame.text:
                return True
        return False

    @staticmethod
    def _remove_shape(shape) -> None:
        shape.element.getparent().remove(shape.element)

    @staticmethod
    def _create_output_presentation(template_prs: Presentation) -> Presentation:
        prs = Presentation()
        prs.slide_width = template_prs.slide_width
        prs.slide_height = template_prs.slide_height
        return prs

    @staticmethod
    def _style_cell_text(cell, bold: bool, size: int) -> None:
        para = cell.text_frame.paragraphs[0]
        if not para.runs:
            para.add_run()
        para.runs[0].font.bold = bold
        para.runs[0].font.size = Pt(size)
        para.alignment = PP_ALIGN.CENTER

    def _resolve_yoy_color(self, metric_key: str, yoy_value) -> RGBColor:
        if yoy_value is None or pd.isna(yoy_value):
            return self.text_secondary

        lower_is_better = {"Cost", "CPC", "CPL"}
        is_positive = yoy_value >= 0
        is_favorable = (not is_positive) if metric_key in lower_is_better else is_positive
        return self.positive if is_favorable else self.negative


def _hex_to_rgb(hex_value: str) -> RGBColor:
    value = hex_value.lstrip("#")
    if len(value) != 6:
        value = "000000"
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))
